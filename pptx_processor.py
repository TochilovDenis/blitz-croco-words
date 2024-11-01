from pptx import Presentation
import os


def get_open_pptx_with_pptx_library(extract_dir: str, pptx_file: dict[int, str], number: int,):
    """
    Открывает и выводит содержимое выбранной презентации.

    Args:
    extract_dir (str): Директория с извлеченными файлами презентаций
    pptx_file (dict[int, str]): Словарь с номерами и именами файлов презентаций
    number (int): Номер выбранной презентации
    """
    try:
        pptx: str = pptx_file[number]
        path: str = os.path.join(extract_dir, pptx)
        print(path)
        prs = Presentation(path)
        print("Презентация открыта.")
        print(f"Количество слайдов: {len(prs.slides)}")

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                print(shape.text)

    except Exception as e:
        print(f"Ошибка при открытии презентации: {e}")


def is_not_valid(text: str) -> bool:
    return ' ' in text or '-' in text or ':' in text or 'СУПЕРКРОКО' in text or 'БЛИЦ-КРОКОДИЛ' in text


def get_extract_words_from_pptx(extract_dir: str, pptx_file: dict[int, str], number: int) -> list[str]:
    """
    Извлекает слова из презентации.

    Args:
    extract_dir (str): Директория с извлеченными файлами презентаций
    pptx_file (dict[int, str]): Словарь с номерами и именами файлов презентаций
    number (int): Номер выбранной презентации

    Returns:
    list[str]: Список извлеченных слов
    """
    words: list = []
    try:
        pptx: str = pptx_file[number]
        path: str = os.path.join(extract_dir, pptx)
        prs: Presentation = Presentation(path)

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text: str = shape.text
                    if is_not_valid(text):
                        continue
                    words.append(text)
        return words

    except Exception as e:
        print(f"Ошибка при открытии презентации: {e}")




