from zipfile import ZipFile
from pathlib import Path
from pptx import Presentation
from get_current_dir import get_current_dir

def get_list_archive_files(folder: str, zip_file: str) -> None:
    """
    Список всех файлов в zip-архиве.
    Returns:
    list[str]: Список имен файлов в архиве.
    """
    # Открываем zip-файл в текущей директории, в папке который указываем, с именем архива
    with ZipFile(Path(get_current_dir()) / folder / zip_file, 'r') as zFile:
        for file in zFile.namelist():
            print(file)


def get_extract_pptx_from_zip(folder: str, zip_file: str, extraction_dir: str):
    """
    Извлекает файлы презентаций из zip-архива.
    Args:
    src (str): Папка с исходным кодом
    filename_zip (str): Имя файла архива
    extraction_dir (str): Директория для извлечения файлов
    """
    with ZipFile(Path(get_current_dir()) / folder / zip_file, 'r') as zFile:
        zFile.extractall(path=extraction_dir)

    print(f"Файлы из {zip_file} распакованы в {extraction_dir}")


def get_extract_pptx_text(zip_file_paths: Path):
    with ZipFile(Path(get_current_dir())/zip_file_paths, 'r') as zip_ref:
        file_list = zip_ref.namelist()
        pptx_files = [f for f in file_list]

        words = []
        for i, pptx_file in enumerate(pptx_files):
            print(f"\nОбработка файла {i + 1}: {pptx_file}")
            prs = Presentation(zip_ref.open(pptx_file))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text: str = shape.text
                        if ' ' in text or ':' in text or 'БЛИЦ-КРОКОДИЛ' in text:
                            continue
                        words.append(text)
                        print(text)
        return words