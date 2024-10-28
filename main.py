import os  # Импорт модуля os для работы с операционной системой
from zipfile import ZipFile  # Импорт класса ZipFile из модуля zipfile для работы с zip-файлами
from pathlib import Path # Импорт класса Path из модуля pathlib для удобной работы с путями
from pptx import Presentation # Импорт класса Presentation из модуля pptx для работы с презентациями
from random import randint

# Константы
SRC = 'src'  # Папка с исходным кодом
FILENAME_ZIP = 'croco-blitz-source.zip'  # Имя файла архива
EXTRACTION_DIR: str = "PPTX"  # Директория для извлечения файлов презентаций
FILENAME_TXT: str = 'words.txt'


# Словарь с номерами и соответствующими им именами файлов презентаций
PPTX: dict[int, str] = {
    1: "Osennyaya_igra_3.pptx",
    2: "Zimnyaya_igra_1.pptx",
    3: "Osennyaya_igra_12.pptx",
    4: "Osennyaya_igra_11.pptx",
    5: "Osennyaya_igra_10.pptx",
    6: "Osennyaya_igra_9.pptx",
    7: "Osennyaya_igra_6.pptx",
    8: "Osennyaya_igra_5.pptx",
    9: "Osennyaya_igra_4.pptx"
}


def current_dir() -> str:
    """
    Получает путь к текущей директории.
    Returns:
    str: Путь к текущей директории.
    """

    # Получаем текущий каталог, в котором находится скрипт
    current_folder: str = os.path.dirname(os.path.realpath(__file__))

    # Возвращаем путь к текущему каталогу
    return current_folder


def list_archive_files(folder: str, filename_zip: str) -> None:
    """
    Список всех файлов в zip-архиве.
    Returns:
    list[str]: Список имен файлов в архиве.
    """
    # Открываем zip-файл в текущей директории, в папке src, с именем FILENAME_ZIP
    with ZipFile(Path(current_dir()) / folder / filename_zip, 'r') as zFile:
        for file in zFile.namelist():
            print(file)


def extract_pptx_from_zip(src: str, filename_zip: str, extraction_dir: str):
    """
    Извлекает файлы презентаций из zip-архива.
    Args:
    src (str): Папка с исходным кодом
    filename_zip (str): Имя файла архива
    extraction_dir (str): Директория для извлечения файлов
    """
    with ZipFile(Path(current_dir()) / src / filename_zip, 'r') as zFile:
        zFile.extractall(path=extraction_dir)

    print(f"Файлы из {filename_zip} распакованы в {extraction_dir}")


def open_pptx_with_pptx_library(extract_dir: str, pptx_file: dict[int, str], number: int,):
    """
    Открывает и выводит содержимое выбранной презентации.

    Args:
    extract_dir (str): Директория с извлеченными файлами презентаций
    pptx_file (dict[int, str]): Словарь с номерами и именами файлов презентаций
    number (int): Номер выбранной презентации
    """
    try:
        pptx: str = pptx_file[number]    # Получаем имя файла презентации по указанному номеру
        path: str = os.path.join(extract_dir, pptx)  # Создаем полный путь к файлу презентации
        print(path)    # Выводим путь к файлу презентации
        prs = Presentation(path)  # Открываем презентацию
        print("Презентация открыта.")
        print(f"Количество слайдов: {len(prs.slides)}") # Выводим количество слайдов в презентации
        # Проходим по всем слайдам в презентации
        for slide in prs.slides:
            # Для каждого слайда проходим по всем шаблонам (shapes)
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                # print(shape.text_frame)
                print(shape.text)

                # Проверяем, имеет ли шаблон текст
                # if hasattr(shape, 'text'):
                #     # Если у шаблона есть текст, выводим его
                #     print(shape.text)
    except Exception as e:
        print(f"Ошибка при открытии презентации: {e}")


def extract_words_from_pptx(extract_dir: str, pptx_file: dict[int, str], number: int):
    words: list = []
    try:
        pptx: str = pptx_file[number]
        path: str = os.path.join(extract_dir, pptx)
        prs: PPTX = Presentation(path)

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text: str = shape.text
                    if ' ' in text or ':' in text or 'БЛИЦ-КРОКОДИЛ' in text:
                        continue
                    words.append(text)

    except Exception as e:
        print(f"Ошибка при открытии презентации: {e}")

    return words


def write_txt(filename:str, extracted_words):
    with open(filename, 'w', encoding='utf-8') as f:
        for word in extracted_words:
            f.write(word + '\n')
    print(f'Слова сохранены в {filename}')


def main() -> None:
    # Выводим путь к текущему каталогу
    print(f"Путь: {current_dir()}")

    # Выводим список файлов в архиве
    print("Список файлов в архиве: ")
    list_archive_files(SRC, FILENAME_ZIP)
    # Извлекаем файл презентации
    extract_pptx_from_zip(SRC, FILENAME_ZIP, EXTRACTION_DIR)
    # Открываем и выводим содержимое выбранной презентации
    open_file_pptx: int = randint(1, 9)
    open_pptx_with_pptx_library(EXTRACTION_DIR, PPTX, open_file_pptx)

    extracted_words: list[str] = extract_words_from_pptx(EXTRACTION_DIR, PPTX, open_file_pptx)
    write_txt(FILENAME_TXT, extracted_words)


if __name__ == "__main__":
    main()
