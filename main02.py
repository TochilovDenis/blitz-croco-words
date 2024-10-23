import os
from pathlib import Path
from zipfile import ZipFile
from pptx import Presentation

BASE_DIR: str = 'src'
FILENAME_ZIP: str = "croco-blitz-source.zip"


def extract_and_open_presentation():
    """
    Функция извлечения презентации из ZIP-файла и ее открытия с помощью библиотеки pptx.
    """
    current_folder = os.path.dirname(os.path.realpath(__file__))

    # Извлечение презентации из ZIP-архива
    with ZipFile(Path(current_folder) / BASE_DIR / FILENAME_ZIP) as archive:
        archive.extractall(path=current_folder)

        # Находим файл презентации в извлеченных файлах
        presentation_file = (f for f in os.listdir(current_folder) if f.endswith('.pptx'))

        if presentation_file:
            try:
                # Открываем презентацию с помощью объекта Presentation
                prs = Presentation(presentation_file)

                # Выводим информацию о презентации
                print(f"\nСодержимое презентации '{presentation_file}':")

                # Выводим количество слайдов
                print(f"Количество слайдов: {len(prs.slides)}")

                # Выводим количество элементов на первом слайде
                if prs.slides[0].shapes:
                    print(f"Количество элементов на первом слайде: {len(prs.slides[0].shapes)}")

                # Выводим количество слайдов с текстом
                text_slides = sum(
                    1 for slide in prs.slides if any(hasattr(shape, "text_frame") for shape in slide.shapes))
                print(f"Количество слайдов с текстом: {text_slides}")

                # Выводим содержимое первого слайда
                if prs.slides[0]:
                    print("\nСодержимое первого слайда:")
                    for shape in prs.slides[0].shapes:
                        if hasattr(shape, "text_frame"):
                            print(shape.text_frame.text)

                # Выводим содержимое всех слайдов с текстом
                slides_with_text = []
                for slide in prs.slides:
                    if any(hasattr(shape, "text_frame") for shape in slide.shapes):
                        slides_with_text.append(slide)

                for i, slide in enumerate(slides_with_text, start=1):
                    print(f"\nСлайд {i}:")
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame"):
                            print(shape.text_frame.text)

            except Exception as e:
                print(f"\nОшибка при открытии презентации: {str(e)}")
        else:
            print("Файл презентации не найден в ZIP-файле.")


if __name__ == "__main__":
    extract_and_open_presentation()
