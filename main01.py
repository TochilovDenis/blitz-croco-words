import os
from pathlib import Path
from zipfile import ZipFile
from pptx import Presentation

BASE_DIR: str = 'src'
FILENAME_ZIP: str = "croco-blitz-source.zip"


def read_zip_file():
    """
    Функция чтения ZIP-файла и извлечения презентации.
    Просматривает содержимое ZIP-файла, извлекает презентацию и выводит ее содержимое.
    """
    current_folder = os.path.dirname(os.path.realpath(__file__))

    # Извлечение презентации из ZIP-архива
    with ZipFile(Path(current_folder) / BASE_DIR / FILENAME_ZIP) as archive:
        archive.extractall(path=current_folder)

        # Находим презентацию в извлеченных файлах
        presentation_files = [f for f in os.listdir(current_folder) if f.endswith('.pptx')]

        if presentation_files:
            # Читаем презентацию
            prs = Presentation(presentation_files[0])

            # Выводим содержимое презентации
            print("Содержимое презентации:")
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        print(shape.text)
        else:
            print("Презентация не найдена в ZIP-файле.")


if __name__ == "__main__":
    read_zip_file()