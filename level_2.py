from typing import IO, TextIO
from pptx import Presentation
from datetime import datetime

def is_not_valid(text: str) -> bool:
    return ' ' in text or '-' in text or ':' in text or 'СУПЕРКРОКО' in text


def get_words_form_file(file: IO[bytes] | TextIO) -> list[str]:
    print(f"Получение слов из файла {file.name}")

    prs = Presentation(file)

    result: list[str] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text # Извлекаем текст слайда

            if is_not_valid(text):
                continue

            # Получаем дату создания или изменения презентации
            presentation_date: datetime.date = prs.core_properties.modified.date()
            formatted_word: str = f"{text}:{presentation_date}"
            result.append(formatted_word)

    print(f"Получение слов выполнено. получено {len(result)} слов")

    return result


def save_words_to_file(words: list[str] or set[str], filename: str) -> None:
    print(f"Сохранение {len(words)} слов в {filename} файл")

    with open(file=filename, mode='w', encoding='utf-8') as file:
        file.writelines(f"{word}\n" for word in words)

    print(f"Готово")